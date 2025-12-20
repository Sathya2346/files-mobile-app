from django.shortcuts import render

# Create your views here.
from rest_framework.response import Response
from rest_framework.decorators import api_view
from django.contrib.auth.hashers import make_password



# myapp/views.py
from rest_framework.decorators import api_view
from rest_framework.response import Response
from rest_framework import status
from django.contrib.auth.hashers import check_password
from .models import User
from .serializers import UserSerializer

# SIGNUP
from rest_framework.decorators import api_view
from rest_framework.response import Response
from rest_framework import status
from django.contrib.auth.hashers import make_password, check_password
from .models import User


# -----------------------------
# ✔ SIGNUP API
# -----------------------------
@api_view(["POST"])
def signup(request):
    try:
        name = request.data.get("name")
        email = request.data.get("email")
        phone = request.data.get("phone")
        password = request.data.get("password")

        if not all([name, email, phone, password]):
            return Response({"error": "All fields required"}, status=400)

        if User.objects.filter(email=email).exists():
            return Response({"error": "Email already registered"}, status=400)

        user = User.objects.create(
            name=name,
            email=email,
            phone=phone,
            password=make_password(password)
        )

        return Response({
            "message": "User created successfully",
            "userId": user.id,
            "name": user.name,
            "email": user.email,
            "phone": user.phone,
            "token": f"token_{user.id}"
        }, status=201)

    except Exception as e:
        return Response({"error": str(e)}, status=500)



# -----------------------------
# ✔ LOGIN API
# -----------------------------
@api_view(["POST"])
def login(request):
    try:
        email = request.data.get("email")
        password = request.data.get("password")

        if not email or not password:
            return Response({"error": "Email and password required"}, status=400)

        try:
            user = User.objects.get(email=email)
        except User.DoesNotExist:
            return Response({"error": "Invalid email"}, status=404)

        if not check_password(password, user.password):
            return Response({"error": "Wrong password"}, status=401)

        return Response({
            "message": "Login successful",
            "userId": user.id,
            "email": user.email,
            "name": user.name,
            "phone": user.phone,
            "token": f"token_{user.id}"
        }, status=200)

    except Exception as e:
        return Response({"error": str(e)}, status=500)



# -----------------------------
# ✔ PROFILE API
# -----------------------------
@api_view(["GET"])
def get_user_profile(request, user_id):
    try:
        user = User.objects.get(id=user_id)

        return Response({
            "id": user.id,
            "name": user.name,
            "email": user.email,
            "phone": user.phone,
            "bio": user.bio,
            "location": user.location,
            "photo": user.photo,
            "created_at": user.created_at
        }, status=200)

    except User.DoesNotExist:
        return Response({"error": "User not found"}, status=404)

    except Exception as e:
        return Response({"error": str(e)}, status=500)



@api_view(['PUT'])
def update_user_profile(request, user_id):
    """Update user profile"""
    try:
        print(f"📥 PUT Profile Request - User ID: {user_id}")
        print(f"📥 Request Data: {request.data}")
        
        user = User.objects.get(id=user_id)
        
        # Update allowed fields
        if "name" in request.data:
            user.name = request.data["name"]
        if "email" in request.data:
            # Check if email already exists for another user
            if User.objects.filter(email=request.data["email"]).exclude(id=user_id).exists():
                return Response({"error": "Email already exists"}, status=status.HTTP_400_BAD_REQUEST)
            user.email = request.data["email"]
        if "phone" in request.data:
            user.phone = request.data["phone"]
        if "bio" in request.data:
            user.bio = request.data["bio"]
        if "location" in request.data:
            user.location = request.data["location"]
        if "photo" in request.data:
            user.photo = request.data["photo"]
        
        # Save without re-hashing password
        user.save(update_fields=['name', 'email', 'phone', 'bio', 'location', 'photo'])
        
        response_data = {
            "id": user.id,
            "name": user.name,
            "email": user.email,
            "phone": user.phone,
            "bio": getattr(user, 'bio', ''),
            "location": getattr(user, 'location', ''),
            "photo": getattr(user, 'photo', ''),
            "created_at": user.created_at.isoformat()
        }
        
        print(f"✅ Profile Updated: {response_data}")
        return Response(response_data, status=status.HTTP_200_OK)
        
    except User.DoesNotExist:
        print(f"❌ User not found: {user_id}")
        return Response({"error": "User not found"}, status=status.HTTP_404_NOT_FOUND)
    except Exception as e:
        print(f"❌ Error updating profile: {str(e)}")
        print(traceback.format_exc())
        return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    


@api_view(['GET'])
def welcome(request):
    return Response({"message": "Welcome from Django Backend!"})



import io
from django.http import JsonResponse, FileResponse
from django.views.decorators.csrf import csrf_exempt
from PyPDF2 import PdfMerger

@csrf_exempt
def merge_pdfs(request):
    if request.method == "POST":
        files = request.FILES.getlist("files")
        if not files or len(files) < 2:
            return JsonResponse({"error": "Please upload at least two PDF files"}, status=400)

        merger = PdfMerger()
        for f in files:
            merger.append(f)

        output = io.BytesIO()
        merger.write(output)
        merger.close()
        output.seek(0)

        # ✅ Return merged PDF as a download (no saving)
        response = FileResponse(output, as_attachment=True, filename="merged.pdf")
        response["Content-Type"] = "application/pdf"
        return response

    return JsonResponse({"error": "Invalid request"}, status=400)



import os
import io
import fitz  # PyMuPDF
import tempfile
import threading
import time
from PIL import Image
from django.http import JsonResponse, FileResponse
from django.views.decorators.csrf import csrf_exempt


def delete_later(*paths, delay=10):
    """Delete temp files after delay."""
    def _cleanup():
        time.sleep(delay)
        for p in paths:
            try:
                os.remove(p)
            except:
                pass
    threading.Thread(target=_cleanup, daemon=True).start()

@csrf_exempt
def compress_pdf(request):
    if request.method != "POST":
        return JsonResponse({"error": "Invalid request"}, status=400)

    upload = request.FILES.get("file")
    if not upload:
        return JsonResponse({"error": "No file uploaded"}, status=400)

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_in:
            for chunk in upload.chunks():
                temp_in.write(chunk)
            input_path = temp_in.name

        output_path = input_path.replace(".pdf", "_compressed.pdf")

        import fitz  # PyMuPDF
        from PIL import Image
        import io

        src_doc = fitz.open(input_path)
        dst_doc = fitz.open()

        for page_index in range(src_doc.page_count):
            src_page = src_doc.load_page(page_index)
            pix = src_page.get_pixmap(matrix=fitz.Matrix(1, 1))

            # Convert to Pillow image
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            img_io = io.BytesIO()
            img.save(img_io, format="JPEG", quality=60, optimize=True)
            img_bytes = img_io.getvalue()

            rect = fitz.Rect(0, 0, pix.width, pix.height)
            new_page = dst_doc.new_page(width=rect.width, height=rect.height)
            new_page.insert_image(rect, stream=img_bytes)

            img_io.close()
            img.close()

        dst_doc.save(output_path, deflate=True, garbage=4, clean=True)
        src_doc.close()
        dst_doc.close()

        response = FileResponse(
            open(output_path, "rb"),
            as_attachment=True,
            filename=f"compressed_{upload.name}"
        )

        delete_later(input_path, output_path)
        return response

    except Exception as e:
        print("❌ Compression Error:", e)
        return JsonResponse({"error": "Compression failed", "details": str(e)}, status=500)





# import io
# from django.http import JsonResponse, FileResponse
# from PyPDF2 import PdfReader, PdfWriter
# from django.views.decorators.csrf import csrf_exempt

# @csrf_exempt
# def split_pdf(request):
#     if request.method != "POST":
#         return JsonResponse({"error": "Invalid request"}, status=400)

#     upload = request.FILES.get("file")
#     page_range = request.POST.get("range", "").strip()

#     if not upload or not page_range:
#         return JsonResponse({"error": "File and range are required"}, status=400)

#     # Validate page range format
#     try:
#         start, end = [int(x) for x in page_range.split("-")]
#         if start < 1 or end < start:
#             raise ValueError()
#     except Exception:
#         return JsonResponse({"error": "Invalid range format (use 1-3)"}, status=400)

#     try:
#         reader = PdfReader(upload)
#         writer = PdfWriter()

#         total_pages = len(reader.pages)
#         end = min(end, total_pages)

#         for i in range(start - 1, end):
#             writer.add_page(reader.pages[i])

#         output = io.BytesIO()
#         writer.write(output)
#         output.seek(0)

#         filename = f"split_{upload.name}"

#         return FileResponse(
#             output,
#             as_attachment=True,
#             filename=filename,
#             content_type="application/pdf"
#         )

#     except Exception as e:
#         return JsonResponse({"error": f"Failed to split PDF: {str(e)}"}, status=500)

import io
import zipfile
from django.http import JsonResponse, HttpResponse
from PyPDF2 import PdfReader, PdfWriter
from django.views.decorators.csrf import csrf_exempt


@csrf_exempt
def split_pdf(request):
    if request.method != "POST":
        return JsonResponse({"error": "Invalid request"}, status=400)

    upload = request.FILES.get("file")
    ranges = request.POST.get("range", "").strip()

    if not upload or not ranges:
        return JsonResponse({"error": "File and range are required"}, status=400)

    # Parse multiple ranges: "1-2,4-6,8-9"
    try:
        range_list = []
        for r in ranges.split(","):
            r = r.strip()
            if "-" not in r:
                raise ValueError()

            start, end = [int(x) for x in r.split("-")]
            if start < 1 or end < start:
                raise ValueError()

            range_list.append((start, end))

    except Exception:
        return JsonResponse({"error": "Invalid range format (use 1-3 or 1-2,4-6)"}, status=400)

    try:
        reader = PdfReader(upload)
        total_pages = len(reader.pages)

        # Create ZIP to store multiple PDFs
        zip_buffer = io.BytesIO()
        zip_file = zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED)

        part_number = 1

        for start, end in range_list:
            end = min(end, total_pages)

            writer = PdfWriter()
            for i in range(start - 1, end):
                writer.add_page(reader.pages[i])

            pdf_bytes = io.BytesIO()
            writer.write(pdf_bytes)
            pdf_bytes.seek(0)

            # Add each split PDF inside ZIP
            zip_file.writestr(f"split_part_{part_number}.pdf", pdf_bytes.read())
            part_number += 1

        zip_file.close()
        zip_buffer.seek(0)

        response = HttpResponse(zip_buffer, content_type="application/zip")
        response["Content-Disposition"] = 'attachment; filename="split_results.zip"'
        return response

    except Exception as e:
        return JsonResponse({"error": f"Failed to split PDF: {str(e)}"}, status=500)



from django.http import JsonResponse, FileResponse
from django.views.decorators.csrf import csrf_exempt
from pdf2docx import Converter
import tempfile, os

@csrf_exempt
def pdf_to_word(request):
    if request.method == 'POST':
        uploaded_file = request.FILES.get('file')  # ✅ must be 'file'
        if not uploaded_file:
            return JsonResponse({'error': "No file uploaded. Use form field name 'file'."}, status=400)

        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            for chunk in uploaded_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name

        temp_docx_path = temp_pdf_path.replace('.pdf', '.docx')

        try:
            cv = Converter(temp_pdf_path)
            cv.convert(temp_docx_path)
            cv.close()

            response = FileResponse(open(temp_docx_path, 'rb'), as_attachment=True, filename='converted.docx')
            return response

        except Exception as e:
            return JsonResponse({'error': str(e)}, status=500)

        finally:
            try:
                os.remove(temp_pdf_path)
                if os.path.exists(temp_docx_path):
                    os.remove(temp_docx_path)
            except:
                pass

    return JsonResponse({'error': 'Invalid request method'}, status=405)





# myapp/views.py
# import os
# import tempfile
# import subprocess
# from django.http import JsonResponse, FileResponse
# from django.views.decorators.csrf import csrf_exempt

# def delete_later(*paths, delay=5):
#     import threading, time
#     def _cleanup():
#         time.sleep(delay)
#         for p in paths:
#             try:
#                 os.remove(p)
#             except:
#                 pass
#     threading.Thread(target=_cleanup, daemon=True).start()

# @csrf_exempt
# def word_to_pdf(request):
#     """
#     Accepts form-data 'file' (.docx/.doc). Returns converted .pdf as stream.
#     Uses LibreOffice (soffice) in headless mode to convert.
#     """
#     if request.method != "POST":
#         return JsonResponse({"error": "POST required"}, status=405)

#     uploaded = request.FILES.get("file")
#     if not uploaded:
#         return JsonResponse({"error": "No file uploaded. Use form field name 'file'."}, status=400)

#     # Accept docx/doc
#     fname = uploaded.name.lower()
#     if not (fname.endswith(".docx") or fname.endswith(".doc")):
#         return JsonResponse({"error": "Upload a .docx or .doc file."}, status=400)

#     try:
#         # write temp input file
#         with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded.name)[1]) as tmp_in:
#             for chunk in uploaded.chunks():
#                 tmp_in.write(chunk)
#             tmp_in.flush()
#             tmp_in_path = tmp_in.name

#         # create temp output path (LibreOffice will create file next to input)
#         tmp_out_dir = tempfile.mkdtemp()
#         # call soffice to convert
#         # --headless --convert-to pdf --outdir <tmp_out_dir> <tmp_in_path>
#         cmd = ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmp_out_dir, tmp_in_path]
#         proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=120)
#         if proc.returncode != 0:
#             err = proc.stderr.decode(errors="ignore") or proc.stdout.decode(errors="ignore")
#             # cleanup
#             try:
#                 os.remove(tmp_in_path)
#             except:
#                 pass
#             return JsonResponse({"error": "Conversion failed", "details": err}, status=500)

#         base = os.path.splitext(os.path.basename(tmp_in_path))[0]
#         pdf_path = os.path.join(tmp_out_dir, base + ".pdf")
#         if not os.path.exists(pdf_path):
#             # cleanup
#             try:
#                 os.remove(tmp_in_path)
#             except:
#                 pass
#             return JsonResponse({"error": "Converted file not found"}, status=500)

#         # Stream back PDF
#         response = FileResponse(open(pdf_path, "rb"), as_attachment=True, filename=base + ".pdf", content_type="application/pdf")

#         # schedule cleanup of temp files and directory
#         delete_later(tmp_in_path, pdf_path)
#         # remove tmp_out_dir after a short delay
#         try:
#             os.rmdir(tmp_out_dir)
#         except:
#             pass

#         return response

#     except subprocess.TimeoutExpired:
#         try:
#             os.remove(tmp_in_path)
#         except:
#             pass
#         return JsonResponse({"error": "Conversion timed out"}, status=500)
#     except Exception as e:
#         # best-effort cleanup
#         try:
#             os.remove(tmp_in_path)
#         except:
#             pass
#         return JsonResponse({"error": "Conversion error", "details": str(e)}, status=500)

import os, tempfile, subprocess, platform
from django.http import JsonResponse, FileResponse
from django.views.decorators.csrf import csrf_exempt

@csrf_exempt
def word_to_pdf(request):
    if request.method != "POST" or "file" not in request.FILES:
        return JsonResponse({"error": "No file uploaded"}, status=400)

    word_file = request.FILES["file"]
    ext = ".docx" if word_file.name.lower().endswith(".docx") else ".doc"

    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp_in:
        tmp_in.write(word_file.read())
        tmp_in.flush()
        tmp_in_path = tmp_in.name

    out_pdf_path = tmp_in_path.rsplit(".", 1)[0] + ".pdf"

    try:
        if platform.system() == "Windows":
            libreoffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
            if not os.path.exists(libreoffice_path):
                return JsonResponse({"error": "LibreOffice not found"}, status=500)

            result = subprocess.run(
                [libreoffice_path, "--headless", "--convert-to", "pdf",
                 "--outdir", os.path.dirname(tmp_in_path), tmp_in_path],
                capture_output=True, text=True
            )
            if result.returncode != 0:
                raise Exception(result.stderr or "LibreOffice conversion failed")
        else:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf",
                 "--outdir", os.path.dirname(tmp_in_path), tmp_in_path],
                check=True
            )

        if not os.path.exists(out_pdf_path):
            raise FileNotFoundError("PDF not generated")

        response = FileResponse(open(out_pdf_path, "rb"), content_type="application/pdf")
        response["Content-Disposition"] = f'attachment; filename="{os.path.basename(out_pdf_path)}"'
        return response

    except Exception as e:
        return JsonResponse({"error": "Conversion error", "details": str(e)}, status=500)

    finally:
        for f in [tmp_in_path, out_pdf_path]:
            try:
                if os.path.exists(f):
                    os.remove(f)
            except:
                pass






import tempfile, os
from django.http import JsonResponse, FileResponse
from django.views.decorators.csrf import csrf_exempt
from PIL import Image

@csrf_exempt
def image_to_pdf(request):
    if request.method != "POST" or "images" not in request.FILES:
        return JsonResponse({"error": "No images uploaded"}, status=400)

    images = request.FILES.getlist("images")
    if len(images) == 0:
        return JsonResponse({"error": "No images uploaded"}, status=400)

    try:
        temp_files = []
        pil_images = []

        for img_file in images:
            temp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
            temp.write(img_file.read())
            temp.flush()
            temp_files.append(temp.name)
            pil_img = Image.open(temp.name).convert("RGB")
            pil_images.append(pil_img)

        pdf_path = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf").name

        # Save images as PDF (first image + rest as append)
        pil_images[0].save(pdf_path, save_all=True, append_images=pil_images[1:])

        response = FileResponse(open(pdf_path, "rb"), content_type="application/pdf")
        response["Content-Disposition"] = 'attachment; filename="converted.pdf"'
        return response

    except Exception as e:
        return JsonResponse({"error": "Conversion failed", "details": str(e)}, status=500)

    finally:
        for f in temp_files + [pdf_path]:
            try:
                if os.path.exists(f):
                    os.remove(f)
            except:
                pass



# import base64
# from django.http import JsonResponse
# from django.views.decorators.csrf import csrf_exempt
# from pdf2image import convert_from_bytes
# from io import BytesIO

# @csrf_exempt
# def pdf_to_images(request):
#     if request.method != "POST":
#         return JsonResponse({"error": "Only POST allowed"}, status=405)

#     pdf_file = request.FILES.get("pdf")
#     if not pdf_file:
#         return JsonResponse({"error": "No PDF uploaded"}, status=400)

#     try:
#         # Convert PDF to images
#         images = convert_from_bytes(pdf_file.read(), dpi=150)
#         result = []

#         for img in images:
#             buffer = BytesIO()
#             img.save(buffer, format="PNG")
#             b64 = base64.b64encode(buffer.getvalue()).decode("utf-8")
#             result.append(f"data:image/png;base64,{b64}")

#         return JsonResponse({"images": result}, status=200)

#     except Exception as e:
#         return JsonResponse({"error": "Conversion failed", "details": str(e)}, status=500)
# views.py
# views.py
import json
import base64
import zipfile
from io import BytesIO
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from pdf2image import convert_from_bytes


@csrf_exempt
def pdf_to_images(request):
    """
    Receives PDF file.
    Returns JSON:
    {
        "images": ["data:image/png;base64,...", ...]
    }
    """
    if request.method != "POST":
        return JsonResponse({"error": "Only POST allowed"}, status=405)

    pdf_file = request.FILES.get("pdf")
    if not pdf_file:
        return JsonResponse({"error": "No PDF uploaded"}, status=400)

    try:
        pages = convert_from_bytes(pdf_file.read(), dpi=150)
        result = []

        for img in pages:
            buf = BytesIO()
            img.save(buf, format="PNG")
            b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
            result.append(f"data:image/png;base64,{b64}")

        return JsonResponse({"images": result}, status=200)

    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@csrf_exempt
def zip_selected_images(request):
    """
    Receives JSON:
    {
        "images": ["data:image/png;base64,...", ...],
        "indexes": [0,2,3]
    }
    Returns JSON:
    {
        "zip_base64": "<base64 zip data>"
    }
    """
    if request.method != "POST":
        return JsonResponse({"error": "Only POST allowed"}, status=405)

    try:
        data = json.loads(request.body)
        images = data.get("images", [])
        indexes = data.get("indexes", [])

        if not images or not indexes:
            return JsonResponse({"error": "Missing images or indexes"}, status=400)

        # ZIP creation
        zip_buf = BytesIO()
        with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
            for idx in indexes:
                if 0 <= idx < len(images):
                    base64_img = images[idx].split(",")[1]
                    img_bytes = base64.b64decode(base64_img)
                    zf.writestr(f"page_{idx+1}.png", img_bytes)

        zip_buf.seek(0)
        zip_b64 = base64.b64encode(zip_buf.getvalue()).decode("utf-8")
        return JsonResponse({"zip_base64": zip_b64}, status=200)

    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)



# backend/views.py
import base64
from io import BytesIO

from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from PIL import Image, ImageOps


@csrf_exempt
def compress_image(request):
    """
    POST: form-data with key 'image' (file),
          optional form fields: max_width (int), quality (int 1-95)
    Returns JSON:
      {
        "image": "data:image/jpeg;base64,...",
        "orig_bytes": 12345,
        "compressed_bytes": 4567,
        "width": 800,
        "height": 600,
        "format": "JPEG"
      }
    All processing in-memory; nothing written to disk.
    """
    if request.method != "POST":
        return JsonResponse({"error": "Only POST allowed"}, status=405)

    img_file = request.FILES.get("image")
    if not img_file:
        return JsonResponse({"error": "No image uploaded"}, status=400)

    try:
        # Read optional params
        try:
            max_width = int(request.POST.get("max_width")) if request.POST.get("max_width") else None
        except Exception:
            max_width = None

        try:
            quality = int(request.POST.get("quality")) if request.POST.get("quality") else 75
        except Exception:
            quality = 75

        quality = max(5, min(95, quality))  # clamp

        # Open image in memory
        img = Image.open(img_file)
        img = ImageOps.exif_transpose(img)  # respect orientation

        orig_buffer = BytesIO()
        img.save(orig_buffer, format="PNG")  # store original size (PNG lossless)
        orig_bytes = orig_buffer.tell()

        # Resize if needed (maintain aspect ratio)
        if max_width and isinstance(max_width, int) and img.width > max_width:
            ratio = max_width / float(img.width)
            new_height = int(img.height * ratio)
            img = img.resize((max_width, new_height), Image.LANCZOS)

        # Convert to RGB if necessary (JPEG needs RGB)
        if img.mode in ("RGBA", "LA") or (img.mode == "P"):
            background = Image.new("RGB", img.size, (255, 255, 255))
            background.paste(img, mask=getattr(img, "split", lambda: [None])[3] if img.mode == "RGBA" else None)
            img = background
        else:
            img = img.convert("RGB")

        out_buffer = BytesIO()
        # Save as JPEG with requested quality
        img.save(out_buffer, format="JPEG", quality=quality, optimize=True)
        compressed_bytes = out_buffer.tell()
        out_buffer.seek(0)
        b64 = base64.b64encode(out_buffer.read()).decode("utf-8")
        data_url = f"data:image/jpeg;base64,{b64}"

        return JsonResponse(
            {
                "image": data_url,
                "orig_bytes": orig_bytes,
                "compressed_bytes": compressed_bytes,
                "width": img.width,
                "height": img.height,
                "format": "JPEG",
                "quality_used": quality,
            },
            status=200,
        )
    except Exception as e:
        return JsonResponse({"error": "Compression failed", "details": str(e)}, status=500)




# # myapp/views.py
# from django.http import FileResponse
# from PyPDF2 import PdfReader, PdfWriter
# import os
# from django.conf import settings
# from django.views.decorators.csrf import csrf_exempt

# @csrf_exempt  # <<< Exempt CSRF for mobile requests
# def protect_pdf(request):
#     if request.method != "POST":
#         return FileResponse(status=405)

#     pdf_file = request.FILES.get("pdf")
#     password = request.POST.get("password", "1234")

#     if not pdf_file:
#         return FileResponse(status=400)

#     input_path = os.path.join(settings.MEDIA_ROOT, "input.pdf")
#     output_path = os.path.join(settings.MEDIA_ROOT, "protected.pdf")

#     # Save uploaded PDF
#     with open(input_path, "wb") as f:
#         for chunk in pdf_file.chunks():
#             f.write(chunk)

#     # Encrypt PDF
#     reader = PdfReader(input_path)
#     writer = PdfWriter()
#     for page in reader.pages:
#         writer.add_page(page)

#     writer.encrypt(user_pwd=password, owner_pwd=password, use_128bit=True)

#     with open(output_path, "wb") as f:
#         writer.write(f)

#     return FileResponse(open(output_path, "rb"), as_attachment=True, filename="protected.pdf")

# myapp/views.py
from django.http import HttpResponse
from PyPDF2 import PdfReader, PdfWriter
from io import BytesIO
from django.views.decorators.csrf import csrf_exempt

@csrf_exempt
def protect_pdf(request):
    if request.method != "POST":
        return HttpResponse("Method not allowed", status=405)

    pdf_file = request.FILES.get("pdf")
    password = request.POST.get("password", "1234")

    if not pdf_file:
        return HttpResponse("No PDF uploaded", status=400)

    # Read PDF into memory
    pdf_bytes = pdf_file.read()
    reader = PdfReader(BytesIO(pdf_bytes))
    writer = PdfWriter()

    # Copy pages
    for page in reader.pages:
        writer.add_page(page)

    # Encrypt
    writer.encrypt(user_pwd=password, owner_pwd=password, use_128bit=True)

    # Save output in memory
    output_stream = BytesIO()
    writer.write(output_stream)
    output_stream.seek(0)

    # Send response
    response = HttpResponse(output_stream, content_type="application/pdf")
    response['Content-Disposition'] = 'attachment; filename="protected.pdf"'
    return response




# myapp/views.py
from django.http import HttpResponse, JsonResponse
from django.views.decorators.csrf import csrf_exempt
from PyPDF2 import PdfReader, PdfWriter
from io import BytesIO

@csrf_exempt
def unlock_pdf(request):
    """
    Accepts POST multipart/form-data with:
      - 'pdf' file
      - 'password' string

    Returns unlocked PDF as an attachment, processed entirely in memory.
    """
    if request.method != "POST":
        return JsonResponse({"detail": "Method not allowed"}, status=405)

    pdf_file = request.FILES.get("pdf")
    password = request.POST.get("password", "")

    if not pdf_file or password is None:
        return JsonResponse({"detail": "PDF and password required"}, status=400)

    try:
        # Read uploaded file into memory
        input_bytes = pdf_file.read()
        input_stream = BytesIO(input_bytes)

        # Load PDF reader
        reader = PdfReader(input_stream)

        # If it's encrypted, try decrypting
        if reader.is_encrypted:
            # PyPDF2's decrypt returns 0 if fail, >0 if success
            try:
                decrypt_result = reader.decrypt(password)
            except Exception:
                decrypt_result = 0

            if decrypt_result == 0:
                return JsonResponse({"detail": "Incorrect password"}, status=400)

        # Create writer and copy pages
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)

        # Output to memory (unlocked PDF)
        output_stream = BytesIO()
        writer.write(output_stream)
        output_stream.seek(0)
        response = HttpResponse(output_stream.read(), content_type="application/pdf")
        response["Content-Disposition"] = 'attachment; filename="unlocked.pdf"'
        return response

    except Exception as e:
        # don't leak internals in production; send minimal info
        return JsonResponse({"detail": "Failed to unlock PDF", "error": str(e)}, status=500)


# myapp/views.py
import io
import tempfile
from django.http import FileResponse, JsonResponse
from django.views.decorators.csrf import csrf_exempt
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

@csrf_exempt
def pdf_to_ppt(request):
    if request.method == "POST":
        pdf_file = request.FILES.get("file")
        if not pdf_file:
            return JsonResponse({"error": "No file uploaded"}, status=400)

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_pdf:
            tmp_pdf.write(pdf_file.read())
            tmp_pdf.flush()

            # Convert PDF pages to images
            images = convert_from_path(tmp_pdf.name)

            # Create a PowerPoint presentation
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]

            for img in images:
                slide = prs.slides.add_slide(blank_slide_layout)
                temp_img = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
                img.save(temp_img.name, "JPEG")
                slide.shapes.add_picture(temp_img.name, Inches(0), Inches(0), Inches(10), Inches(7.5))

            ppt_stream = io.BytesIO()
            prs.save(ppt_stream)
            ppt_stream.seek(0)

            return FileResponse(
                ppt_stream,
                as_attachment=True,
                filename="converted.pptx",
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
    return JsonResponse({"error": "Invalid request"}, status=400)





import os
import subprocess
import platform
from django.http import HttpResponse
from rest_framework.decorators import api_view
from tempfile import mkdtemp

@api_view(["POST"])
def ppt_to_pdf(request):
    try:
        file = request.FILES.get("file")
        if not file:
            return HttpResponse("No file uploaded", status=400)

        temp_dir = mkdtemp()
        input_path = os.path.join(temp_dir, file.name)
        output_path = os.path.join(temp_dir, os.path.splitext(file.name)[0] + ".pdf")

        # Save uploaded PPT/PPTX
        with open(input_path, "wb") as f:
            for chunk in file.chunks():
                f.write(chunk)

        # Detect system path for LibreOffice
        if platform.system() == "Windows":
            soffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
        else:
            soffice_path = "soffice"  # Linux & macOS

        # Convert PPT → PDF (headless)
        subprocess.run([
            soffice_path,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", temp_dir,
            input_path,
        ], check=True)

        # Read and return converted PDF
        with open(output_path, "rb") as f:
            pdf_data = f.read()

        response = HttpResponse(pdf_data, content_type="application/pdf")
        response["Content-Disposition"] = 'attachment; filename=\"converted.pdf\"'
        return response

    except Exception as e:
        print("PPT→PDF error:", e)
        return HttpResponse(str(e), status=500)







# # myapp/views.py
import io
import fitz  # PyMuPDF
from django.http import HttpResponse
from rest_framework.decorators import api_view

@api_view(["POST"])
def add_watermark(request):
    try:
        file = request.FILES.get("file")
        if not file:
            return HttpResponse("No file uploaded", status=400)

        text = request.POST.get("text", "CONFIDENTIAL")
        position = request.POST.get("position", "middle")  # top, middle, bottom
        align = request.POST.get("align", "horizontal")    # horizontal or vertical
        rotate = float(request.POST.get("rotate", "0"))
        opacity = float(request.POST.get("opacity", "0.2"))

        # Load PDF into memory
        pdf_bytes = file.read()
        pdf_in = fitz.open(stream=pdf_bytes, filetype="pdf")

        font_size = 40
        text_color = (0, 0, 0)  # black

        for page in pdf_in:
            width, height = page.rect.width, page.rect.height

            # Determine Y position
            if position == "top":
                y = 100
            elif position == "bottom":
                y = height - 100
            else:  # middle
                y = height / 2

            x = width / 2  # center X

            # Adjust rotation for vertical
            text_rotate = rotate
            if align == "vertical":
                text_rotate += 90  # add 90 degrees for vertical alignment

            # Insert watermark at point to avoid wrapping
            # PyMuPDF uses point = top-left corner
            # For vertical, adjust x/y to roughly center
            if align == "vertical":
                page.insert_text(
                    fitz.Point(x, y - font_size/2),
                    text,
                    fontsize=font_size,
                    rotate=text_rotate,
                    color=text_color,
                    fill_opacity=opacity,
                    render_mode=0,  # normal fill
                )
            else:
                # Horizontal: center by shifting x by half text width
                # Approximate width = font_size * len(text) * 0.5
                x_shift = x - (font_size * len(text) * 0.25)
                page.insert_text(
                    fitz.Point(x_shift, y),
                    text,
                    fontsize=font_size,
                    rotate=text_rotate,
                    color=text_color,
                    fill_opacity=opacity,
                    render_mode=0,
                )

        # Save output to memory
        output_stream = io.BytesIO()
        pdf_in.save(output_stream)
        pdf_in.close()

        response = HttpResponse(output_stream.getvalue(), content_type="application/pdf")
        response["Content-Disposition"] = 'attachment; filename="watermarked.pdf"'
        return response

    except Exception as e:
        print("Watermark Error:", e)
        return HttpResponse(f"Error: {e}", status=500)

# import io
# import fitz  # PyMuPDF
# from django.http import HttpResponse
# from rest_framework.decorators import api_view

# @api_view(["POST"])
# def add_watermark(request):
#     try:
#         file = request.FILES.get("file")
#         if not file:
#             return HttpResponse("No file uploaded", status=400)

#         text = request.POST.get("text", "CONFIDENTIAL")
#         position = request.POST.get("position", "middle")  # top, middle, bottom
#         align = request.POST.get("align", "horizontal")    # horizontal or vertical
#         opacity = float(request.POST.get("opacity", "0.2"))

#         pdf_bytes = file.read()
#         pdf_in = fitz.open(stream=pdf_bytes, filetype="pdf")

#         font_size = 40
#         text_color = (0, 0, 0)  # black

#         for page in pdf_in:
#             width, height = page.rect.width, page.rect.height

#             if align == "vertical":
#                 # Vertical: top to bottom
#                 x = width / 2  # center horizontally
#                 if position == "top":
#                     start_y = 100
#                 elif position == "bottom":
#                     start_y = height - font_size * len(text) - 100
#                 else:  # middle
#                     start_y = (height - font_size * len(text)) / 2

#                 for i, char in enumerate(text):
#                     page.insert_text(
#                         fitz.Point(x - font_size/2, start_y + i * font_size),
#                         char,
#                         fontsize=font_size,
#                         color=text_color,
#                         fill_opacity=opacity,
#                         render_mode=0,
#                     )

#             else:
#                 # Horizontal: top/middle/bottom
#                 if position == "top":
#                     y = 100
#                 elif position == "bottom":
#                     y = height - 100
#                 else:  # middle
#                     y = height / 2

#                 # Approx center horizontally
#                 x_shift = (width - font_size * len(text) * 0.5) / 2
#                 page.insert_text(
#                     fitz.Point(x_shift, y),
#                     text,
#                     fontsize=font_size,
#                     color=text_color,
#                     fill_opacity=opacity,
#                     render_mode=0,
#                 )

#         output_stream = io.BytesIO()
#         pdf_in.save(output_stream)
#         pdf_in.close()

#         response = HttpResponse(output_stream.getvalue(), content_type="application/pdf")
#         response["Content-Disposition"] = 'attachment; filename="watermarked.pdf"'
#         return response

#     except Exception as e:
#         print("Watermark Error:", e)
#         return HttpResponse(f"Error: {e}", status=500)



import os
import io
import uuid
import camelot
import pdfplumber
from openpyxl import Workbook
from django.http import HttpResponse
from rest_framework.decorators import api_view

@api_view(["POST"])
def pdf_to_excel(request):
    try:
        file = request.FILES.get("file")
        if not file:
            return HttpResponse("No file uploaded", status=400)

        # Save uploaded PDF temporarily
        pdf_bytes = file.read()
        temp_pdf_path = f"temp_{uuid.uuid4().hex}.pdf"
        with open(temp_pdf_path, "wb") as f:
            f.write(pdf_bytes)

        excel_wb = Workbook()
        ws = excel_wb.active
        ws.title = "Sheet1"

        # Try Camelot (tables)
        tables = camelot.read_pdf(temp_pdf_path, pages="all", flavor="stream")
        row_count = 1
        if tables.n > 0:
            for table in tables:
                for row in table.data:
                    for c, cell in enumerate(row):
                        ws.cell(row=row_count, column=c+1, value=cell)
                    row_count += 1
        else:
            # Fallback: pdfplumber (text extraction)
            with pdfplumber.open(temp_pdf_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        for line in text.split("\n"):
                            ws.cell(row=row_count, column=1, value=line)
                            row_count += 1

        # Save Excel to memory
        output_stream = io.BytesIO()
        excel_wb.save(output_stream)
        output_stream.seek(0)

        # Remove temp PDF
        os.remove(temp_pdf_path)

        return HttpResponse(
            output_stream.getvalue(),
            content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

    except Exception as e:
        print("PDF→Excel Error:", e)
        return HttpResponse(f"Error: {e}", status=500)










# # myapp/views.py
# import io
# from openpyxl import load_workbook
# from reportlab.lib.pagesizes import A4
# from reportlab.pdfgen import canvas
# from django.http import HttpResponse
# from rest_framework.decorators import api_view

# @api_view(["POST"])
# def excel_to_pdf(request):
#     try:
#         file = request.FILES.get("file")
#         if not file:
#             return HttpResponse("No file uploaded", status=400)

#         # Load Excel workbook
#         excel_bytes = file.read()
#         workbook = load_workbook(io.BytesIO(excel_bytes))
#         sheet = workbook.active

#         # Create PDF in memory
#         output_stream = io.BytesIO()
#         pdf = canvas.Canvas(output_stream, pagesize=A4)
#         width, height = A4

#         # Set starting positions
#         x_start = 50
#         y_start = height - 50
#         row_height = 20

#         # Loop through Excel rows
#         for r, row in enumerate(sheet.iter_rows(values_only=True)):
#             y = y_start - r * row_height
#             for c, cell in enumerate(row):
#                 x = x_start + c * 100  # cell width approx 100
#                 pdf.drawString(x, y, str(cell) if cell is not None else "")
#             # If y goes below margin, create new page
#             if y - row_height < 50:
#                 pdf.showPage()
#                 y_start = height - 50

#         pdf.save()
#         output_stream.seek(0)

#         response = HttpResponse(
#             output_stream.getvalue(),
#             content_type="application/pdf"
#         )
#         response["Content-Disposition"] = 'attachment; filename="converted.pdf"'
#         return response

#     except Exception as e:
#         print("Excel→PDF Error:", e)
#         return HttpResponse(f"Error: {e}", status=500)
import io
from openpyxl import load_workbook
from reportlab.lib.pagesizes import A4, landscape
from reportlab.pdfgen import canvas
from rest_framework.decorators import api_view
from django.http import HttpResponse
from datetime import datetime, date

@api_view(["POST"])
def excel_to_pdf(request):
    try:
        file = request.FILES.get("file")
        if not file:
            return HttpResponse("No file uploaded", status=400)

        # Load Excel workbook in memory
        wb = load_workbook(file)
        ws = wb.active

        # PDF buffer in memory
        buffer = io.BytesIO()
        c = canvas.Canvas(buffer, pagesize=landscape(A4))
        width, height = landscape(A4)
        margin = 30
        y_start = height - margin
        row_height = 20

        # Determine number of columns per page
        max_col_width = 100  # adjust cell width in points
        cols_per_page = int((width - 2*margin) // max_col_width)
        total_cols = ws.max_column
        pages_needed = (total_cols + cols_per_page - 1) // cols_per_page

        for page in range(pages_needed):
            x = margin
            y = y_start
            start_col = page * cols_per_page + 1
            end_col = min(start_col + cols_per_page - 1, total_cols)

            # Draw rows
            for row in ws.iter_rows(min_row=1, max_row=ws.max_row, min_col=start_col, max_col=end_col, values_only=True):
                x = margin
                for cell in row:
                    # ✅ Format dates correctly
                    if isinstance(cell, (datetime, date)):
                        cell_text = cell.strftime("%Y-%m-%d")  # only date
                    else:
                        cell_text = str(cell) if cell is not None else ""

                    # 🔹 Draw the formatted text
                    c.drawString(x, y, cell_text)
                    x += max_col_width
                y -= row_height
                if y < margin:
                    c.showPage()
                    y = y_start

            c.showPage()  # new page for next set of columns

        c.save()
        buffer.seek(0)
        return HttpResponse(
            buffer.getvalue(),
            content_type="application/pdf"
        )

    except Exception as e:
        print("Excel→PDF Error:", e)
        return HttpResponse(f"Error: {e}", status=500)





# # myapp/views.py
# import io
# from django.http import HttpResponse
# from rest_framework.decorators import api_view
# from PyPDF2 import PdfReader, PdfWriter
# import base64
# from PIL import Image
# import fitz  # PyMuPDF

# @api_view(["POST"])
# def sign_pdf(request):
#     """
#     Accepts PDF file + signature image (base64) + position.
#     Returns signed PDF in memory.
#     """
#     try:
#         pdf_file = request.FILES.get("file")
#         signature_base64 = request.data.get("signature")
#         x = float(request.data.get("x", 100))
#         y = float(request.data.get("y", 100))
#         page_number = int(request.data.get("page", 0))

#         if not pdf_file or not signature_base64:
#             return HttpResponse("Missing PDF or signature", status=400)

#         # Read PDF bytes once
#         pdf_bytes = pdf_file.read()
#         pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")

#         if page_number >= len(pdf_doc):
#             return HttpResponse("Invalid page number", status=400)

#         page = pdf_doc[page_number]

#         # Decode base64 signature
#         if "," in signature_base64:
#             signature_base64 = signature_base64.split(",")[1]  # remove data:image/png;base64,
#         sig_bytes = base64.b64decode(signature_base64)
#         sig_image = Image.open(io.BytesIO(sig_bytes))

#         # Optional: scale signature to reasonable size (e.g., 150x50)
#         max_width = 150
#         max_height = 50
#         w_ratio = max_width / sig_image.width
#         h_ratio = max_height / sig_image.height
#         scale = min(1, w_ratio, h_ratio)
#         sig_width = sig_image.width * scale
#         sig_height = sig_image.height * scale

#         # Insert image
#         rect = fitz.Rect(x, y, x + sig_width, y + sig_height)
#         page.insert_image(rect, stream=sig_bytes)

#         # Save to memory
#         output_stream = io.BytesIO()
#         pdf_doc.save(output_stream)
#         output_stream.seek(0)

#         return HttpResponse(
#             output_stream.getvalue(),
#             content_type="application/pdf"
#         )

#     except Exception as e:
#         print("Sign PDF Error:", e)
#         return HttpResponse(f"Error: {e}", status=500)





# myapp/views.py
import io
import base64
from io import BytesIO
from django.http import HttpResponse
from rest_framework.decorators import api_view
from rest_framework.response import Response
from PIL import Image
import fitz  # PyMuPDF
import os

TEMP_DIR = "/tmp/pdf_uploads"
os.makedirs(TEMP_DIR, exist_ok=True)

# ------------------------
# 📌 PDF → IMAGE PREVIEW
# ------------------------
@api_view(["POST"])
def pdf_to_image(request):
    try:
        pdf = request.FILES.get("file")
        if not pdf:
            return Response({"error": "No PDF provided"}, status=400)

        # Save to temp file
        temp_path = f"{TEMP_DIR}/preview_{pdf.name}"
        with open(temp_path, "wb") as f:
            for chunk in pdf.chunks():
                f.write(chunk)

        # Render first page to PNG at dpi=160 (same dpi client expects)
        doc = fitz.open(temp_path)
        page = doc.load_page(0)

        pix = page.get_pixmap(dpi=160)  # same dpi used in client logic
        # get raw PNG bytes
        img_bytes = pix.pil_tobytes(format="PNG")

        base64_img = "data:image/png;base64," + base64.b64encode(img_bytes).decode()

        return Response({
            "preview": base64_img,
            "width": pix.width,
            "height": pix.height
        })

    except Exception as e:
        print("PDF→Image ERROR:", e)
        return Response({"error": str(e)}, status=500)

@api_view(["POST"])
def sign_pdf(request):
    try:
        pdf = request.FILES.get("file")
        sig_base64 = request.data.get("signature")
        x_top = float(request.data.get("x"))
        y_top = float(request.data.get("y"))   # top-origin from RN
        scale = float(request.data.get("scale", 1))

        if not pdf or not sig_base64:
            return Response({"error": "Missing data"}, status=400)

        pdf_path = f"{TEMP_DIR}/sign_{pdf.name}"
        with open(pdf_path, "wb") as f:
            for chunk in pdf.chunks():
                f.write(chunk)

        if "," in sig_base64:
            _, b64data = sig_base64.split(",", 1)
        else:
            b64data = sig_base64

        sig_bytes = base64.b64decode(b64data)
        sig_img = Image.open(BytesIO(sig_bytes)).convert("RGBA")
        sig_w_px, sig_h_px = sig_img.size

        doc = fitz.open(pdf_path)
        page = doc.load_page(0)
        page_rect = page.rect

        pix = page.get_pixmap(dpi=160)
        factor = page_rect.width / pix.width  # convert preview px → PDF points

        # 🔹 Do NOT flip Y
        x_pdf = x_top * factor
        y_pdf = y_top * factor

        sig_w_pdf = sig_w_px * factor * scale
        sig_h_pdf = sig_h_px * factor * scale

        rect = fitz.Rect(
            x_pdf,
            y_pdf,
            x_pdf + sig_w_pdf,
            y_pdf + sig_h_pdf,
        )

        page.insert_image(rect, stream=sig_bytes)

        output = io.BytesIO()
        doc.save(output)
        doc.close()
        output.seek(0)
        return HttpResponse(output.getvalue(), content_type="application/pdf")

    except Exception as e:
        print("SIGN PDF ERROR:", e)
        return Response({"error": str(e)}, status=500)









# myapp/views.py
from rest_framework.decorators import api_view
from rest_framework.response import Response
from rest_framework import status
from django.core.mail import send_mail
from django.conf import settings
from .models import FAQ, ContactInfo, SupportTicket, User
import traceback

# ... (keep your existing login, signup, profile views here)


@api_view(['GET'])
def get_faqs(request):
    """Get all active FAQs"""
    try:
        faqs = FAQ.objects.filter(is_active=True)
        data = [
            {
                "id": faq.id,
                "question": faq.question,
                "answer": faq.answer,
                "order": faq.order
            }
            for faq in faqs
        ]
        return Response(data, status=status.HTTP_200_OK)
    except Exception as e:
        print(f"❌ Error getting FAQs: {str(e)}")
        return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['GET'])
def get_contact_info(request):
    """Get contact information"""
    try:
        contact = ContactInfo.objects.filter(is_active=True).first()
        if contact:
            data = {
                "phone": contact.phone,
                "email": contact.email,
                "whatsapp": contact.whatsapp
            }
            return Response(data, status=status.HTTP_200_OK)
        return Response({"error": "Contact info not found"}, status=status.HTTP_404_NOT_FOUND)
    except Exception as e:
        print(f"❌ Error getting contact info: {str(e)}")
        return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['POST'])
def submit_ticket(request):
    """Submit a support ticket and send email to admin"""
    try:
        print(f"📥 Submit Ticket Request: {request.data}")
        
        user_id = request.data.get('userId')
        issue_type = request.data.get('issueType')
        message = request.data.get('message')
        user_email = request.data.get('userEmail')
        user_name = request.data.get('userName')
        
        if not all([user_id, issue_type, message]):
            return Response(
                {"error": "Missing required fields"},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Get user
        try:
            user = User.objects.get(id=user_id)
        except User.DoesNotExist:
            return Response(
                {"error": "User not found"},
                status=status.HTTP_404_NOT_FOUND
            )
        
        # Create ticket
        ticket = SupportTicket.objects.create(
            user=user,
            issue_type=issue_type,
            message=message
        )
        
        # Send email to admin
        try:
            admin_email = settings.ADMIN_EMAIL  # Set this in settings.py
            subject = f"New Support Ticket #{ticket.id} - {issue_type.upper()}"
            
            email_message = f"""
New Support Ticket Received

Ticket ID: #{ticket.id}
User: {user_name} ({user_email})
Issue Type: {issue_type.upper()}
Status: OPEN

Message:
{message}

---
Submitted on: {ticket.created_at.strftime('%Y-%m-%d %H:%M:%S')}

Please login to admin panel to respond: {settings.ADMIN_PANEL_URL}
            """
            
            send_mail(
                subject=subject,
                message=email_message,
                from_email=settings.DEFAULT_FROM_EMAIL,
                recipient_list=[admin_email],
                fail_silently=False,
            )
            
            print(f"✅ Email sent to admin: {admin_email}")
            
        except Exception as email_error:
            print(f"⚠️ Email sending failed: {str(email_error)}")
            # Don't fail the request if email fails
        
        print(f"✅ Ticket created: #{ticket.id}")
        
        return Response({
            "message": "Ticket submitted successfully",
            "ticketId": ticket.id
        }, status=status.HTTP_201_CREATED)
        
    except Exception as e:
        print(f"❌ Error submitting ticket: {str(e)}")
        print(traceback.format_exc())
        return Response(
            {"error": str(e)},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )